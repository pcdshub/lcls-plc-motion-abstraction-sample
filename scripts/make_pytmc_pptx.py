import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title, content, image_path=None, code=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.5)
    width = prs.slide_width - Inches(1)
    height = prs.slide_height - Inches(2.4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    
    # Content text
    p = frame.add_paragraph()
    p.text = content
    p.font.size = Pt(20)
    
    # Code block (example)
    if code:
        p = frame.add_paragraph()
        p.text = code
        p.font.size = Pt(14)
        p.font.name = 'Courier New'
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    if image_path:
        slide.shapes.add_picture(image_path, Inches(6.0), Inches(2.2), width=Inches(3.2), height=Inches(2.0))

# ---- Create visuals ----
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

def make_flow_diagram(filename):
    fig, ax = plt.subplots(figsize=(7, 5))
    ax.axis('off')

    # Box specs: (text, (x, y), color)
    boxes = [
        ("TwinCAT\nProject File\n(.tsproj/.tmc)", (0.1, 0.7), '#FFD700'),  # gold
        ("XML Parsing\n(lxml)", (0.35, 0.7), '#B0E0E6'),  # light blue
        ("Recursive\nElement Walk", (0.6, 0.7), '#9ACD32'),  # yellowgreen
        ("Class Name Resolution\n(Special for <Symbol>)", (0.85, 0.7), '#AFEEEE'),  # paleturquoise
        ("Dynamic Class Creation\nSymbol_ST_MotionStage", (0.85, 0.4), '#ADD8E6'),  # light blue
        ("Object Tree Built", (0.6, 0.4), '#90EE90'),  # lightgreen
        ("Search for Symbol\n(find(Symbol_ST_MotionStage))", (0.35, 0.4), '#FFB6C1'),  # light pink
    ]
    for text, (x, y), color in boxes:
        ax.text(x, y, text, ha='center', va='center', fontsize=12,
                bbox=dict(boxstyle="round,pad=0.6", fc=color, ec='black', lw=2))

    # Arrows (from left-to-right, then down and looping back)
    arrowprops = dict(arrowstyle="->", linewidth=2, color='black')
    ax.annotate('', xy=(0.28,0.7), xytext=(0.18,0.7), arrowprops=arrowprops)
    ax.annotate('', xy=(0.52,0.7), xytext=(0.38,0.7), arrowprops=arrowprops)
    ax.annotate('', xy=(0.76,0.7), xytext=(0.62,0.7), arrowprops=arrowprops)
    
    ax.annotate('', xy=(0.85,0.62), xytext=(0.85,0.53), arrowprops=arrowprops)
    ax.annotate('', xy=(0.80,0.43), xytext=(0.65,0.43), arrowprops=arrowprops)
    ax.annotate('', xy=(0.55,0.43), xytext=(0.40,0.43), arrowprops=arrowprops)
    ax.annotate('', xy=(0.32,0.58), xytext=(0.32,0.45), arrowprops=arrowprops)

    plt.tight_layout()
    plt.savefig(filename, dpi=120)
    plt.close()

def make_symbol_nc_connection(filename):
    fig, ax = plt.subplots(figsize=(7, 3.6))
    ax.axis('off')
    # Symbol box
    ax.text(0.08, 0.5, 'Symbol\nST_MotionStage', ha='center', va='center', fontsize=13,
            bbox=dict(boxstyle="round,pad=0.5", fc="#E5FCFF", ec="#2196F3", lw=2))
    # Link box
    ax.text(0.34, 0.5, 'Link\n(VarA/VarB)', ha='center', va='center', fontsize=13,
            bbox=dict(boxstyle="round,pad=0.5", fc="#FFF2CC", ec="#FF9800", lw=2))
    # NC box
    ax.text(0.60, 0.5, 'NC', ha='center', va='center', fontsize=13,
            bbox=dict(boxstyle="round,pad=0.5", fc="#FFD6E5", ec="#AD1457", lw=2))
    # Axis box
    ax.text(0.85, 0.5, 'Axis\n(M1)', ha='center', va='center', fontsize=13,
            bbox=dict(boxstyle="round,pad=0.5", fc="#D0FFD6", ec="#388E3C", lw=2))

    # Arrows
    arr_opts = dict(arrowstyle="->", lw=2)
    ax.annotate("", xy=(0.22,0.5), xytext=(0.14,0.5), arrowprops=arr_opts)
    ax.annotate("", xy=(0.48,0.5), xytext=(0.40,0.5), arrowprops=arr_opts)
    ax.annotate("", xy=(0.73,0.5), xytext=(0.67,0.5), arrowprops=arr_opts)
    ax.annotate("", xy=(0.07,0.55), xytext=(0.07,0.70), arrowprops=dict(arrowstyle="->"))
    ax.text(0.07, 0.73, ".nc_axis property", ha='center', va='center', fontsize=10, color="#2196F3")
    plt.tight_layout()
    plt.savefig(filename, dpi=120)
    plt.close(fig)
    
def make_xml_python_mapping(filename):
    fig, ax = plt.subplots(figsize=(7, 2))
    ax.axis('off')
    # XML box
    ax.text(0.05, 0.5,
            '<Symbol>\n  <Name>Main.M1</Name>\n  <BaseType>ST_MotionStage</BaseType>\n</Symbol>',
            va='center', ha='left', fontsize=11, fontfamily='monospace',
            bbox=dict(fc='#FFFAE5', ec='orange', lw=2))
    # Python class box
    ax.text(0.55, 0.5,
            '↓',
            va='center', ha='center', fontsize=36, fontweight='bold')
    ax.text(0.7, 0.5,
            'Python object:\nSymbol_ST_MotionStage',
            va='center', ha='left', fontsize=13, fontfamily='monospace',
            bbox=dict(fc='#E5FCFF', ec='#2196F3', lw=2))
    # Optionally add an image/logo
    plt.tight_layout()
    plt.savefig(filename, dpi=120)
    plt.close()
    
# ----- Main PowerPoint Creation -----
prs = Presentation()
prs.slide_width = Inches(12)
prs.slide_height = Inches(7)

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Parsing TwinCAT Project Files to Python Objects"
slide.placeholders[1].text = "How pytmc transforms XML files into rich, type-specific Python symbol classes"

# Slide 2: Overview
add_slide(prs, "Process Overview",
      "- Parses .tsproj or .tmc XML files\n"
      "- Builds an in-memory tree of Python objects\n"
      "- Dynamically creates custom symbol classes\n"
      "- Example: ST_MotionStage → Symbol_ST_MotionStage")

# Slide 3: Parsing Entry Point
add_slide(prs, "Step 1: Parsing Entry Point",
          "User calls parse('my_project.tsproj') to start parsing.\n"
          "Code reads the XML, finds the root, and hands off to the recursive parser.",
          code="project = parse('my_project.tsproj')")

# Slide 4: Recursive XML → Python Mapping
add_slide(prs, "Step 2: Recursive XML → Python Mapping",
          "TwincatItem.parse(element, parent, filename) is called for each XML node.\n"
          " - Decides appropriate Python class for each tag\n"
          " - Recurses into child elements.",
          code="TwincatItem.parse(root, filename, parent=None)")

# Slide 5: Handling Symbol Elements
make_xml_python_mapping("xml_python_mapping.png")
add_slide(prs, "Step 3: Handling Symbol Elements",
    "When parsing a <Symbol> element:\n"
    " - Looks up <BaseType>\n"
    " - Generates a custom Python class: Symbol_{BaseType} (e.g. Symbol_ST_MotionStage)",
    image_path="xml_python_mapping.png")

# Slide 6: Dynamic Class Creation
add_slide(prs, "How Dynamic Class Creation Works",
          "Checks if the class exists in the registry.\n"
          "If not, dynamically creates and registers the class.",
          code="""cls = type("Symbol_ST_MotionStage", (Symbol,), {})
TWINCAT_TYPES["Symbol_ST_MotionStage"] = cls""")

# Slide 7: Building the Object Tree
add_slide(prs, "Step 4: Building the Object Tree",
          "Each XML element becomes a Python object instance.\n"
          "Objects keep references to parents and children.\n"
          "Symbols are available for searching/manipulation.")

# Slide 8: Accessing Special Symbols
add_slide(prs, "Step 5: Accessing Special Symbols",
          "You can now search for special symbol types.",
          code="""motion_stage_symbol = next(project.find(Symbol_ST_MotionStage))
print(motion_stage_symbol)""")

# Slide 9: Why Dynamic Subclassing?
add_slide(prs, "Why Dynamic Subclassing?",
          "- Enables type-specific logic/behavior\n"
          "- Easy searching (find(Symbol_ST_MotionStage))\n"
          "- Extensible for new symbol types")
# Slide: How is Symbol ST_MotionStage Connected to NC Axis?
add_slide(prs, "How is ST_MotionStage Symbol Connected to NC Axis?",
    "Connection is established via TwinCAT's 'Link' objects, using naming conventions and XML references. "
    "\n\nThe Symbol's .nc_axis property:"
    "\n- Finds the relevant Link for the symbol"
    "\n- Resolves the corresponding NC and axis name"
    "\n- Returns the correct Axis object",
    image_path=None)

# Add the diagram showing this connection
make_symbol_nc_connection("symbol_nc_connection.png")
add_slide(prs, "Symbol–NC Axis Connection Visual",
    "ST_MotionStage symbol connects to an NC axis through a Link object."
    "\nSee how these objects relate in the parsed tree:",
    image_path="symbol_nc_connection.png")

# Slide: Example from XML to Python Access
add_slide(prs, "Example: XML and Python Code for Axis Connection",
    "Let's trace the connection from XML definition, through parsing, to Python object lookup.",
    code="""# XML fragments
<NC>
    <Axis name="M1" ... />
</NC>
...
<Symbol>
    <Name>Main.M1</Name>
    <BaseType>ST_MotionStage</BaseType>
</Symbol>
...
<Link VarA="^Main.M1.Axis.NcToPlc" VarB="^TINC.Task1.AxisSection.M1.Axis.NcToPlc" ... />

# Python code
sym = next(project.find(Symbol_ST_MotionStage))
axis = sym.nc_axis    # Returns the Axis object for 'M1'
""")
# Slide 10: Complete Flow Diagram
make_flow_diagram("flow_diagram.png")
add_slide(prs, "Complete Flow Diagram",
          "Full process to parse a TwinCAT project file and create symbol objects.",
          image_path="flow_diagram.png")

# Slide 11: Summary
add_slide(prs, "Summary",
          "- XML files are safely mapped to Python objects\n"
          "- Custom Symbol classes dynamically created\n"
          "- Handles complex TwinCAT project configurations")

# Slide 12: Q&A
add_slide(prs, "Q&A", "Questions?")

prs.save("pytmc_twinCAT_parsing.pptx")
print("Presentation saved as 'pytmc_twinCAT_parsing.pptx'")
