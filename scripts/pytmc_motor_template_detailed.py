from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt

def add_slide(prs, title, bullet_points=None, image=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    if bullet_points is not None:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    if image:
        # Insert image with a suitable size below the text box
        img_left = Inches(1)
        img_top = Inches(2)
        img_width = Inches(8)
        slide.shapes.add_picture(image, img_left, img_top, width=img_width)
    return slide

def add_code_slide(prs, title, code, bg_rgb=(245,245,220)):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.2)
    width = prs.slide_width - Inches(1.0)
    height = prs.slide_height - Inches(2.0)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.text = code
    p.font.size = Pt(11)
    p.font.name = 'Consolas'
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_rgb)
    tf.word_wrap = True

def make_modern_pytmc_diagram(filename="pytmc_motor_diagram.png"):
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.axis('off')
    blocks = [
        ("TwinCAT\n(.tsproj/.tmc)", (0.10, 0.65), 0.19, 0.12, "#B3E5FC"),
        ("pytmc parser/\npragmas/record", (0.36, 0.65), 0.23, 0.12, "#DCEDC8"),
        ("get_motors", (0.62, 0.7), 0.14, 0.06, "#FFD54F"),
        ("Legacy\nST_MotionStage", (0.82, 0.58), 0.15, 0.08, "#FFF9C4"),
        ("FB_MotionStage\n(axis-link)", (0.82, 0.75), 0.15, 0.09, "#FFECB3"),
        ("EthercatMC\ncontroller & axes", (0.62, 0.57), 0.18, 0.09, "#E1BEE7"),
        ("motor.db\n(per axis macros)", (0.62, 0.80), 0.18, 0.09, "#FFE0B2"),
        ("Unified IOC Jinja2\ntemplate", (0.22, 0.37), 0.47, 0.13, "#BBDEFB"),
        ("st.cmd output", (0.5, 0.13), 0.15, 0.08, "#C5CAE9"),
    ]
    for text, (x, y), w, h, color in blocks:
        ax.add_patch(plt.Rectangle((x, y), w, h, fc=color, ec='black', lw=2))
        ax.text(x + w/2, y + h/2, text, ha='center', va='center', fontsize=12)
    ax.annotate("", xy=(0.29, 0.71), xytext=(0.20, 0.71), arrowprops=dict(arrowstyle="->", lw=2))
    ax.annotate("", xy=(0.58, 0.71), xytext=(0.47, 0.71), arrowprops=dict(arrowstyle="->", lw=2))
    ax.annotate("", xy=(0.78, 0.67), xytext=(0.74, 0.67), arrowprops=dict(arrowstyle="->", lw=2))
    ax.annotate("", xy=(0.78, 0.80), xytext=(0.74, 0.725), arrowprops=dict(arrowstyle="->", lw=2))
    ax.annotate("", xy=(0.70, 0.44), xytext=(0.38, 0.44), arrowprops=dict(arrowstyle="->", lw=2))
    ax.annotate("", xy=(0.57, 0.44), xytext=(0.57, 0.69), arrowprops=dict(arrowstyle="->", lw=2))
    ax.annotate("", xy=(0.57, 0.44), xytext=(0.57, 0.85), arrowprops=dict(arrowstyle="->", lw=2))
    ax.annotate("", xy=(0.57, 0.44), xytext=(0.82, 0.85), arrowprops=dict(arrowstyle="->", lw=2))
    ax.annotate("", xy=(0.57, 0.44), xytext=(0.82, 0.62), arrowprops=dict(arrowstyle="->", lw=2))
    ax.annotate("", xy=(0.57, 0.44), xytext=(0.82, 0.79), arrowprops=dict(arrowstyle="->", lw=2))
    ax.annotate("", xy=(0.57, 0.20), xytext=(0.57, 0.44), arrowprops=dict(arrowstyle="->", lw=2))
    plt.tight_layout()
    plt.savefig(filename, bbox_inches='tight', dpi=120)
    plt.close(fig)

prs = Presentation()
prs.slide_width = Inches(12)
prs.slide_height = Inches(7)

prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = "Modernizing pytmc template: Legacy to FB_MotionStage & StreamDevice"

add_slide(prs, "Modernizing pytmc template for Next-Generation Motion", [
    "From legacy ST_MotionStage/EthercatMC to FB_MotionStage/StreamDevice patterns"
])

add_slide(prs, "Motivation for Change", [
    "Beckhoff Baustein trend: axes as FB_MotionStage",
    "Need flexible StreamDevice control and per-axis macros",
    "Clear migration/separation, robust autosave/archiver/OPI handling"
])

add_slide(prs, "Legacy vs. Modern IOC Motor Support", [
    "Legacy: Only EthercatMC (asyn/ai/bo...); Modern: StreamDevice, per-axis macros",
    "Legacy: EthercatMC.template etc.; Modern: motor.db + macro expansion",
    "Legacy: Limited macros, mostly global; Modern: Fully macroized, per axis",
    "Mix: Not supported in IOC, project-level switch only"
])

add_slide(prs, "Motor Symbol Discovery: get_motors Changes", [
    "Legacy: Discovers only Symbol_ST_MotionStage",
    "New: Discovers both ST and FB motors; attaches is_fb_motionstage flag for template filtering",
    "axis-link pragma required for FB motors",
    "Ensures exactly one family per IOC"
])
add_code_slide(prs, "get_motors: Legacy vs Modern", 
"""# Legacy
def get_motors(plc):
    symbols = get_symbols_by_type(plc)
    return [
        stage
        for stage in symbols.get("Symbol_ST_MotionStage", [])
        if not stage.is_pointer and pragmas.has_pragma(stage)
    ]

# Modern
def get_motors(plc):
    symbols = get_symbols_by_type(plc)
    st_motors = ensure_list(symbols.get("Symbol_ST_MotionStage", []))
    fb_motors = ensure_list(symbols.get("Symbol_FB_MotionStage", []))
    motors = []
    for stage in st_motors + fb_motors:
        if not stage.is_pointer and pragmas.has_pragma(stage) and has_axis_link(stage):
            is_fb = (stage.__class__.__name__ == "Symbol_FB_MotionStage")
            setattr(stage, "is_fb_motionstage", is_fb)
            motors.append(stage)
    return motors
""", (255,255,240))

add_slide(prs, "Record Generator Changes: has_axis_link and ensure_list", [
    "has_axis_link ensures FB_MotionStage has axis-link pragma (motor validity gate)",
    "ensure_list ensures all parsing/loops can handle list/nil/single uniformly"
])
add_code_slide(prs, "has_axis_link and ensure_list",
"""def has_axis_link(stage):
    require_axis_link = isinstance(stage, parser.Symbol_FB_MotionStage)
    aliases = getattr(stage, "symbol_aliases", [])
    required_aliases = {'Symbol_FB_MotionStageNC', 'Symbol_FB_MotionStageNCDS402'}
    if set(aliases) & required_aliases:
        require_axis_link = True
    if not require_axis_link:
        return True
    for pragma in pragmas.get_pragma(stage, name="pytmc"):
        for line in pragma.splitlines():
            if line.strip().startswith("axis-link:"):
                return True
    return False

def ensure_list(val):
    if isinstance(val, list):
        return val
    elif val is None:
        return []
    else:
        return [val]
""", (240,255,244))

add_slide(prs, "Template: Legacy EthercatMC", [
    "Legacy motors use classic EthercatMC records and controller/axis instantiation."
])
add_code_slide(prs, "Template: Legacy EthercatMC", 
"""{% if legacy_motors %}
EthercatMCCreateController("$(MOTOR_PORT)", "$(ASYN_PORT)", "$(NUMAXES)", ...)
{% for motor in legacy_motors %}
    # Macros for axis
    EthercatMCCreateAxis("$(MOTOR_PORT)", "$(AXIS_NO)", ...)
    dbLoadRecords("EthercatMC.template", ...)
    ...
{% endfor %}
{% endif %}
""", (255,255,220))

add_slide(prs, "Template: Modern FB_MotionStage with StreamDevice", [
    "Each FB axis loads motor.db with axis macros. No legacy EthercatMC code executes."
])
add_code_slide(prs, "Template: Modern FB_MotionStage with StreamDevice", 
"""{% for motor in fb_motors %}
    epicsEnvSet("MOTOR_NAME", "{{motor|epics_suffix}}")
    dbLoadRecords("motor.db", "PORT=$(ASYN_PORT), ADSPORT=$(AMS_PORT), ADSPATH=$(MOTOR_ADS_PATH), PREFIX=$(MOTOR_PREFIX), M=$(MOTOR_NAME)")
{% endfor %}
""", (224,247,250))

add_slide(prs, "StreamDevice and Bulk Command Support", [
    "Records in motor.db loaded per axis with axis-specific macros",
    "Command/status (and StreamDevice trigger) records are auto-generated"
])
add_code_slide(prs, "StreamDevice Command Example (motor.db)",
"""record(bo, "$(PREFIX):$(M):MOVE_CMD")
{
    field(DESC, "Trigger axis move batch via StreamDevice")
    field(DTYP, "stream")
    field(OUT, "@motor.proto move_cmd($(ADSPORT),$(ADSPATH),$(PREFIX),$(M)) $(PORT)")
}
""", (231,245,233))

add_slide(prs, "Template Logic: Family Selection", [
    "Ensures only one family block is active per IOC",
    "All per-axis macros available regardless of type"
])
add_code_slide(prs, "Final Workflow and dbLoadRecords",
"""# ASYN/trace config always set once
{% if motors %}
    ... macros ...
{% endif %}
# Only one family block runs per IOC
dbLoadRecords("motor.db", ...)  # for FB
EthercatMCCreateController(...) # for legacy
""", (255,242,231))

make_modern_pytmc_diagram("pytmc_motor_diagram.png")
add_slide(prs, "Modern pytmc Template Workflow", image="pytmc_motor_diagram.png")

add_slide(prs, "Summary & Next Steps", [
    "Clean, robust template for EthercatMC or FB_MotionStage motors",
    "All macros and records auto-configured per axis",
    "StreamDevice and command/status records handled with correct macros",
    "Ready for large-scale facility automation"
])

prs.save("pytmc_motor_template_detailed.pptx")
print("Presentation 'pytmc_motor_template_detailed.pptx' and diagram PNG generated!")
